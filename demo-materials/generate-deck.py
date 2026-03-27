"""
Generate slide deck (PPTX) and talk track (PDF) for Code Apps Tech Series.
Advanced concepts + live demo segue. Companion to a colleague's intro talk.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Colors ──────────────────────────────────────────────────────────
PURPLE = RGBColor(0x74, 0x27, 0x74)       # Power Platform purple
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)      # Dark slide background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)

# ── Helpers ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Segoe UI", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Support (text, accent_color) tuples or plain strings
        if isinstance(item, tuple):
            txt, accent = item
            # Bold prefix before " -- "
            if " -- " in txt:
                prefix, rest = txt.split(" -- ", 1)
                run1 = p.add_run()
                run1.text = "\u2022  " + prefix + " "
                run1.font.size = Pt(font_size)
                run1.font.color.rgb = accent
                run1.font.bold = True
                run1.font.name = font_name
                run2 = p.add_run()
                run2.text = rest
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = color
                run2.font.name = font_name
            else:
                run = p.add_run()
                run.text = "\u2022  " + txt
                run.font.size = Pt(font_size)
                run.font.color.rgb = accent
                run.font.name = font_name
        else:
            run = p.add_run()
            run.text = "\u2022  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
    return tf

def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_code_box(slide, left, top, width, height, text, font_size=11):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = "Consolas"
    return tf

# ── PPTX Generation ────────────────────────────────────────────────
def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 1 — TITLE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.08, PURPLE)

    add_text_box(slide, 0.8, 1.5, 11.7, 1.2,
                 "Code Apps: Under the Hood",
                 font_size=44, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 2.8, 11.7, 0.8,
                 "Advanced patterns, runtime internals, and an AI-built dashboard",
                 font_size=22, color=LIGHT_GRAY)
    add_accent_bar(slide, 0.8, 3.8, 3.0, 0.04, ACCENT_VIOLET)
    add_text_box(slide, 0.8, 4.2, 6, 0.6,
                 "Power Platform Tech Series  \u00b7  2026",
                 font_size=16, color=MID_GRAY)

    # Tag line at bottom
    add_text_box(slide, 0.8, 6.2, 11.7, 0.5,
                 "Companion to \"Intro to Code Apps\" \u2014 we pick up where that leaves off",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 2 — WHAT'S ACTUALLY RUNNING
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, PURPLE)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "What's Actually Running",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Your React SPA is a small island inside a massive enterprise runtime",
                 font_size=16, color=MID_GRAY)

    # Left column — the runtime layers
    add_text_box(slide, 0.8, 1.8, 5.5, 0.5,
                 "The Platform Wraps You", font_size=20, color=ACCENT_VIOLET, bold=True)
    add_bullet_list(slide, 0.8, 2.3, 5.5, 4.5, [
        ("AppType: \"ClassicCanvasApp\" -- Your React 19 SPA is classified the same as a drag-and-drop canvas app", ACCENT_RED),
        ("Iframes all the way down -- Shell \u2192 player boot \u2192 iframe to runtime-app.powerapps.com \u2192 your code", ACCENT_AMBER),
        ("77 script bundles registered -- MSAL, legacy ADAL, WinJS, DOMPurify, Teams, Copilot, Floodgate surveys", ACCENT_BLUE),
        ("230KB localization inlined -- NFC, barcode, Cordova, China ICP, Siri Shortcuts\u2026 in every page load", ACCENT_GREEN),
        ("55 feature gates -- Microsoft toggles capabilities per-tenant around your app", MID_GRAY),
    ], font_size=14, spacing=Pt(8))

    # Right column — what you get for free
    add_text_box(slide, 7.2, 1.8, 5.5, 0.5,
                 "What You Get for Free", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, 7.2, 2.3, 5.5, 4.5, [
        ("Server-side auth -- No tokens in HTML. Entra ID authenticates before the page is served", ACCENT_GREEN),
        ("Copilot sidecar slot -- 320px panel, dedicated JS bundle, 185 references in the shell HTML", ACCENT_VIOLET),
        ("Sovereign cloud routing -- .dynamics.cn, .microsoftdynamics.de, .appsplatform.us \u2014 all wired in", ACCENT_BLUE),
        ("Accessibility scaffolding -- ARIA roles, live regions, skip-to-content, progress bar stages", ACCENT_AMBER),
        ("CSP + nonce security -- Every inline script carries a nonce; CSP delivered via HTTP headers", MID_GRAY),
    ], font_size=14, spacing=Pt(8))

    # Bottom callout
    add_accent_bar(slide, 0.8, 6.6, 11.7, 0.04, PURPLE)
    add_text_box(slide, 0.8, 6.7, 11.7, 0.5,
                 "Source: browser dev tools export of the deployed app \u2014 I'll show you in the demo",
                 font_size=13, color=MID_GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 3 — THINGS ARE NOT ALWAYS WHAT THEY APPEAR
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_RED)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Things Are Not Always What They Appear",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Just because it works locally doesn\u2019t mean it\u2019s right \u2014 or that it\u2019ll work deployed",
                 font_size=16, color=MID_GRAY)

    # Left column — Story 1: Data connections
    add_accent_bar(slide, 0.8, 1.85, 0.06, 4.3, ACCENT_AMBER)
    add_text_box(slide, 1.1, 1.8, 5.3, 0.5,
                 "\u26a0  \"Working\" Data Connections", font_size=20, color=ACCENT_AMBER, bold=True)
    add_text_box(slide, 1.1, 2.3, 5.3, 0.4,
                 "Implemented incorrectly \u2014 still appeared to work",
                 font_size=13, color=LIGHT_GRAY)

    add_bullet_list(slide, 1.1, 2.8, 5.3, 3.5, [
        ("Generated types declare parentcustomerid -- but Dataverse OData returns the GUID as _parentcustomerid_value at runtime", LIGHT_GRAY),
        ("Contact form showed \"None\" for linked accounts -- data was actually correct in Dataverse, the UI just couldn\u2019t read it", ACCENT_AMBER),
        ("Saving still worked fine -- writes used different fields than reads, masking the bug", LIGHT_GRAY),
        ("Fix: runtime fallback pattern -- contact.parentcustomerid ?? (contact as Record)._parentcustomerid_value", ACCENT_GREEN),
    ], font_size=13, spacing=Pt(6))

    # Right column — Story 2: Agent iframe
    add_accent_bar(slide, 7.0, 1.85, 0.06, 4.3, ACCENT_RED)
    add_text_box(slide, 7.3, 1.8, 5.3, 0.5,
                 "\u274c  Agent Iframe: Local \u2260 Deployed", font_size=20, color=ACCENT_RED, bold=True)
    add_text_box(slide, 7.3, 2.3, 5.3, 0.4,
                 "Works perfectly on localhost \u2014 fails completely in production",
                 font_size=13, color=LIGHT_GRAY)

    add_bullet_list(slide, 7.3, 2.8, 5.3, 3.5, [
        ("Direct Line + MSAL token exchange -- copied a pattern that worked in a standalone Azure web app", LIGHT_GRAY),
        ("Code Apps use custom auth schemes (paauth, dynamicauth) -- no public API to extract Bearer tokens", ACCENT_RED),
        ("Local dev masked the problem -- Power Platform\u2019s iframe host obscured the token exchange failures", LIGHT_GRAY),
        ("Fix: pivot to Copilot Studio iframe embed -- same environment = session flows natively, zero packages added", ACCENT_GREEN),
    ], font_size=13, spacing=Pt(6))

    # Bottom callout
    add_accent_bar(slide, 0.8, 6.5, 11.7, 0.04, ACCENT_RED)
    add_text_box(slide, 0.8, 6.6, 11.7, 0.6,
                 "Lesson: the Power Platform runtime is not a standard web host \u2014 test deployed, not just local",
                 font_size=15, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 4 — BUILT BY AI AGENTS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_GREEN)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Built 100% by AI Agents",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "18 phases, zero hand-written scaffolding, every decision tracked",
                 font_size=16, color=MID_GRAY)

    # Left — the build story
    add_text_box(slide, 0.8, 1.8, 5.5, 0.5,
                 "The Agentic Build", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, 0.8, 2.3, 5.5, 2.5, [
        ("Claude Code drove every phase -- Scaffolding, CRUD, UI, AI integration, Kanban board", ACCENT_GREEN),
        ("Iterative discovery -- Each phase started with research prompts, then build, then polish", ACCENT_BLUE),
        ("Tracked decisions -- docs/tracked/ has every prompt, gotcha, and fix documented per phase", ACCENT_AMBER),
        ("Real gotchas found by AI -- Polymorphic lookups, $select bug, HashRouter requirement, OData bind syntax", ACCENT_RED),
    ], font_size=14, spacing=Pt(8))

    # Right — the dashboard story
    add_text_box(slide, 7.2, 1.8, 5.5, 0.5,
                 "The Dashboard: Simpler Than You Think", font_size=20, color=ACCENT_VIOLET, bold=True)
    add_bullet_list(slide, 7.2, 2.3, 5.5, 2.5, [
        ("705 lines, zero chart libraries -- Pure CSS + SVG. Donut chart = circle with stroke-dasharray", ACCENT_VIOLET),
        ("One prompt rebuilt it -- \"Drop Chart.js and rebuild with pure CSS/SVG\" \u2192 done", ACCENT_GREEN),
        ("Every viz is clickable -- KPI cards, status rows, priority bars all drill down to filtered tables", ACCENT_BLUE),
        ("Kanban board: 1,117 lines -- 4 columns, cross-entity drag-drop, floating toolbar, priority gradients", ACCENT_AMBER),
    ], font_size=14, spacing=Pt(8))

    # Bottom — the key stats
    add_accent_bar(slide, 0.8, 5.0, 11.7, 0.04, ACCENT_GREEN)

    # Stats row
    stats = [
        ("6", "Dataverse\nEntities"),
        ("18", "Build\nPhases"),
        ("705", "Dashboard\nLines"),
        ("0", "Chart\nLibraries"),
        ("1", "Deploy\nCommand"),
    ]
    stat_width = 2.0
    stat_start = 1.2
    for i, (num, label) in enumerate(stats):
        x = stat_start + i * 2.3
        add_text_box(slide, x, 5.3, stat_width, 0.6,
                     num, font_size=36, color=ACCENT_GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")
        add_text_box(slide, x, 5.9, stat_width, 0.6,
                     label, font_size=12, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 0.8, 6.7, 11.7, 0.5,
                 "npm run build && pac code push  \u2014  that's the entire deploy",
                 font_size=14, color=ACCENT_GREEN, font_name="Consolas",
                 alignment=PP_ALIGN.LEFT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 5 — PATTERNS THAT PAY OFF
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_BLUE)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Patterns That Pay Off",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Standard React patterns \u2014 the platform doesn't fight you",
                 font_size=16, color=MID_GRAY)

    # Four pattern cards as two rows of two
    patterns = [
        ("TanStack Query Cache", ACCENT_BLUE,
         "Ctrl+K command palette searches\nall 6 entities with zero API calls.\nIt reads the query cache directly.",
         "useQueryClient().getQueryData([\"accounts\"])"),
        ("Zustand Store (27 lines)", ACCENT_AMBER,
         "Quick-create pills in the nav open\nthe right form on the right page.\nNo prop drilling, no context nesting.",
         "useQuickCreateStore().open(\"action-items\", { defaultTaskType })"),
        ("dnd-kit + Dataverse Mutations", ACCENT_GREEN,
         "Drag a card to parking lot \u2192\nPATCH tdvsp_pinned = true.\nWithin-column reorder persists to localStorage.",
         "useSortable() + useMutation(updateItem)"),
        ("Tailwind v4 Dark Mode", ACCENT_VIOLET,
         "Class-based dark mode works inside\nthe Power Platform iframe. Media\nqueries don't \u2014 the host controls them.",
         "@custom-variant dark (&:where(.dark, .dark *))"),
    ]

    for i, (title, accent, desc, code) in enumerate(patterns):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 1.8 + row * 2.6

        add_accent_bar(slide, x, y, 0.06, 2.2, accent)
        add_text_box(slide, x + 0.25, y + 0.05, 5.5, 0.4,
                     title, font_size=18, color=accent, bold=True)
        add_text_box(slide, x + 0.25, y + 0.5, 5.5, 1.1,
                     desc, font_size=13, color=LIGHT_GRAY)
        add_code_box(slide, x + 0.25, y + 1.6, 5.5, 0.4,
                     code, font_size=10)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 6 — LET ME SHOW YOU (demo transition)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_AMBER)

    add_text_box(slide, 0.8, 1.0, 11.7, 1.0,
                 "Let Me Show You",
                 font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, 5.5, 2.3, 2.3, 0.04, ACCENT_AMBER)
    add_text_box(slide, 0.8, 2.7, 11.7, 0.6,
                 "Live Demo",
                 font_size=28, color=ACCENT_AMBER,
                 alignment=PP_ALIGN.CENTER)

    # Demo flow bullets — centered
    demo_items = [
        "\u2776  Dashboard \u2014 KPIs, donut chart, drill down into filtered data",
        "\u2777  Kanban Board \u2014 drag to parking lot, watch Dataverse update",
        "\u2778  Ctrl+K \u2014 search all entities instantly from the query cache",
        "\u2779  AI Extraction \u2014 paste meeting notes, watch action items appear",
        "\u277a  Dark mode toggle \u2014 class-based, works inside the iframe",
        "\u277b  Dev tools peek \u2014 see the runtime we just talked about",
        "\u277c  Deploy \u2014 npm run build && pac code push",
    ]
    tf = add_bullet_list(slide, 2.5, 3.5, 8.3, 3.5, demo_items,
                         font_size=16, color=LIGHT_GRAY, spacing=Pt(6))

    pptx_path = os.path.join(OUT_DIR, "code-apps-under-the-hood.pptx")
    prs.save(pptx_path)
    print(f"PPTX saved: {pptx_path}")
    return pptx_path


# ── PDF Talk Track ──────────────────────────────────────────────────
def build_pdf():
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(25, 20, 25)

    # ── Fonts ───────────────────────────────────────────────────────
    pdf.add_font("Segoe", "", "C:/Windows/Fonts/segoeui.ttf", uni=True)
    pdf.add_font("Segoe", "B", "C:/Windows/Fonts/segoeuib.ttf", uni=True)
    pdf.add_font("Segoe", "I", "C:/Windows/Fonts/segoeuii.ttf", uni=True)
    pdf.add_font("Consolas", "", "C:/Windows/Fonts/consola.ttf", uni=True)

    def heading(text, size=18):
        pdf.set_font("Segoe", "B", size)
        pdf.set_text_color(116, 39, 116)  # Purple
        pdf.cell(0, 10, text, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)

    def subheading(text, size=13):
        pdf.set_font("Segoe", "B", size)
        pdf.set_text_color(40, 40, 60)
        pdf.cell(0, 8, text, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(1)

    def body(text, size=10.5):
        pdf.set_font("Segoe", "", size)
        pdf.set_text_color(50, 50, 50)
        pdf.multi_cell(0, 5.5, text)
        pdf.ln(2)

    def stage_direction(text, size=10):
        pdf.set_font("Segoe", "I", size)
        pdf.set_text_color(120, 120, 140)
        pdf.multi_cell(0, 5, text)
        pdf.ln(2)

    def code(text, size=9):
        pdf.set_font("Consolas", "", size)
        pdf.set_text_color(16, 185, 129)  # Green
        pdf.multi_cell(0, 5, text)
        pdf.set_font("Segoe", "", 10.5)
        pdf.set_text_color(50, 50, 50)
        pdf.ln(2)

    def bullet(text, size=10.5):
        pdf.set_font("Segoe", "", size)
        pdf.set_text_color(50, 50, 50)
        x = pdf.get_x()
        pdf.cell(6, 5.5, "\u2022")
        pdf.multi_cell(0, 5.5, text)
        pdf.ln(1)

    def divider():
        pdf.ln(3)
        pdf.set_draw_color(116, 39, 116)
        pdf.set_line_width(0.3)
        y = pdf.get_y()
        pdf.line(25, y, 185, y)
        pdf.ln(5)

    # ════════════════════════════════════════════════════════════════
    # PAGE 1 — TITLE + OVERVIEW
    # ════════════════════════════════════════════════════════════════
    pdf.add_page()
    pdf.set_font("Segoe", "B", 24)
    pdf.set_text_color(116, 39, 116)
    pdf.cell(0, 14, "Code Apps: Under the Hood", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Segoe", "", 13)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 8, "Talk Track & Speaker Notes", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, "Power Platform Tech Series  |  2026", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    body("This talk track accompanies 6 slides covering advanced Code Apps concepts. "
         "It is designed to follow a colleague's \"Intro to Code Apps\" presentation. "
         "Do not repeat intro material (what Code Apps are, how to scaffold, basic CLI). "
         "Instead, go straight to what makes this real and then get into the live demo.")
    pdf.ln(2)

    subheading("Session Flow")
    bullet("Slides: ~10 minutes (6 slides)")
    bullet("Live demo: ~15-20 minutes")
    bullet("Q&A: remaining time")
    pdf.ln(2)

    subheading("Key Messages")
    bullet("Your React app runs inside a massive enterprise runtime you never see -- that runtime handles auth, Copilot, accessibility, sovereign clouds, and telemetry for you.")
    bullet("This entire app was built by AI agents (Claude Code) across 18 phases. Every decision is tracked in docs/tracked/.")
    bullet("The analytics dashboard is 705 lines with zero chart libraries. Pure CSS + SVG. One prompt rebuilt it.")
    bullet("Standard React patterns (TanStack Query, Zustand, dnd-kit, Tailwind) work perfectly inside Power Platform. The platform doesn't fight you.")
    bullet("Just because something works locally doesn't mean it's implemented correctly or that it will work deployed. The Power Platform runtime is not a standard web host.")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 1: Title")
    stage_direction("[Slide on screen. Brief pause. Don't rush this -- let people read it.]")
    body("\"[Colleague] just walked you through what Code Apps are and how to get started. "
         "I'm going to pick up where that leaves off. We're going to look at what's actually "
         "running when you deploy a Code App, how AI agents built this entire demo app, "
         "and then I'll show you the app live. Let's go.\"")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 2: What's Actually Running")
    stage_direction("[Advance slide. Give people 5 seconds to scan the two columns.]")

    body("\"So here's something that surprised me. I saved the HTML from browser dev tools "
         "on my deployed Code App. The file is 81,000 tokens. And my React app? Nowhere in it. "
         "Not a single reference to Vite, React, Tailwind, or any of my components.\"")
    pdf.ln(1)

    body("\"The platform classifies my React 19 SPA as a ClassicCanvasApp. Same AppType as a "
         "drag-and-drop canvas app. It wraps my app in multiple layers of iframes -- the shell "
         "boots a player, the player loads a runtime from runtime-app.powerapps.com, and that "
         "runtime loads my code.\"")
    pdf.ln(1)

    subheading("Left Column -- The Platform Wraps You")
    bullet("77 script bundles registered on CDN -- MSAL for auth, legacy ADAL still there, WinJS, DOMPurify, Teams integration, Copilot JS, Floodgate surveys, Office Browser Feedback. All downloaded before your first React component mounts.")
    bullet("230KB of localization strings inlined in every page -- strings for NFC, barcode scanning, Cordova plugins, China ICP compliance, Siri Shortcuts. Most will never be used in a browser context, but they're there.")
    bullet("55 feature gates -- Boolean flags that Microsoft can toggle per-tenant. CopilotSidecar, EnableMsalV4 (still off), FloodgateSurveys, CSPForceReportViolation. Your app runs inside a runtime that's constantly being reconfigured around you.")
    pdf.ln(1)

    subheading("Right Column -- What You Get for Free")
    bullet("Auth is entirely server-side. No tokens in the HTML. Entra ID authenticates the user before this page is even served. UseClientSideAuth is false.")
    bullet("Copilot sidecar has its own DOM slot (320px), its own JS bundle, teaching bubbles, record picker, summarization, NL-to-filter -- 185 references in the shell HTML alone.")
    bullet("Sovereign cloud routing is baked in -- .dynamics.cn, .microsoftdynamics.de, .appsplatform.us, .dynamics.com.mcas.ms for MCAS. One codebase, every sovereign cloud.")
    bullet("Accessibility: ARIA roles, aria-live regions, skip-to-content, 4-stage progress bar with aria-valuetext. The platform does the hard parts.")
    pdf.ln(1)

    body("\"The takeaway: when you deploy a Code App, you're not hosting a web page. "
         "You're placing your React SPA inside a governed enterprise runtime that handles "
         "auth, telemetry, Copilot, accessibility, and sovereign cloud routing. You just write React.\"")

    stage_direction("[Pause. Let that land. If someone asks about performance overhead, "
                    "acknowledge it's real -- the 77 bundles and 230KB strings are a cost -- "
                    "but it's the same cost Canvas Apps already pay. You're not adding to it.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 3: Things Are Not Always What They Appear")
    stage_direction("[Advance slide. Two stories side by side. This is the 'trust but verify' moment.]")

    body("\"So we've seen what's running under the hood. Now let me tell you about two times "
         "we got burned -- because this platform is not a standard web host, and things that "
         "look right can be very wrong.\"")
    pdf.ln(1)

    subheading("Left Column -- \"Working\" Data Connections")
    body("\"When we built the contact form, we had an account dropdown. You could create a contact, "
         "link it to an account, save it -- everything worked. But when you opened that contact again, "
         "the dropdown showed 'None.' The account link was gone. Except it wasn't gone.\"")
    pdf.ln(1)
    bullet("The data was correct in Dataverse. The account was linked. The problem was that the generated TypeScript types declare a field called parentcustomerid -- that's what you write to. But at runtime, Dataverse OData returns the GUID in a completely different field: _parentcustomerid_value. The generated types don't even include that field.")
    bullet("So writes worked perfectly -- you set parentcustomerid and the link saved. But reads failed silently -- the code tried to read parentcustomerid, got undefined, and showed 'None.'")
    bullet("The fix was a runtime fallback: try the declared field first, then cast to Record<string, unknown> and check _parentcustomerid_value. Same pattern for the account name -- you can't rely on parentcustomeridname from the SDK either. We had to fetch all accounts and build a manual lookup map.")
    pdf.ln(1)
    body("\"The lesson: your generated types are a contract for writes, not for reads. "
         "The OData response has its own ideas about field names.\"")
    pdf.ln(1)

    subheading("Right Column -- Agent Iframe: Local Does Not Equal Deployed")
    body("\"Second story. We wanted to embed a Copilot Studio agent in the app. I had a pattern "
         "that worked great in another project -- a standalone Azure Static Web App. Direct Line "
         "connection, MSAL token exchange, botframework-webchat. Clean, proven, production-tested.\"")
    pdf.ln(1)
    bullet("We copied the pattern into the Code App. It looked good locally. But Code Apps don't use standard OAuth Bearer tokens. The Power Apps SDK uses custom auth schemes -- paauth, dynamicauth -- with no public API to extract a standard Bearer token. The Direct Line token exchange literally cannot work.")
    bullet("Local dev masked the problem. When you run pac code run locally, you're still inside a Power Platform iframe, but the auth context is different enough that you don't immediately see the failure. You'd only discover this after deploying with pac code push.")
    bullet("The fix was a complete pivot. Instead of Direct Line + botframework-webchat, we embedded the Copilot Studio hosted webchat URL in a simple iframe. Since both the Code App and the agent are in the same Power Platform environment, the user's session flows natively. Zero packages needed. Zero token exchange. It just works.")
    pdf.ln(1)
    body("\"The lesson: don't assume patterns from standalone web apps will transfer into Code Apps. "
         "The auth model is fundamentally different. Always test deployed, not just local.\"")

    stage_direction("[This slide often generates questions. Common ones: "
                    "'Can you use MSAL at all?' (No, not in the traditional sense -- the platform owns the auth context.) "
                    "'Does the Copilot iframe work in all browsers?' (Yes, same session cookies.) "
                    "Keep answers brief -- you'll show the working agent in the demo.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 4: Built 100% by AI Agents")
    stage_direction("[Advance slide. The stats row at the bottom will draw eyes -- let people see the numbers.]")

    body("\"This is the part people always ask about. This entire app -- six Dataverse entities, "
         "full CRUD, analytics dashboard, Kanban board with drag-and-drop, command palette, "
         "AI-powered action item extraction -- was built entirely by Claude Code across 18 phases.\"")
    pdf.ln(1)

    subheading("Left Column -- The Agentic Build")
    bullet("Every phase started with a research prompt. 'How do polymorphic lookups work in Code Apps?' The agent would explore, find the gotchas, then build.")
    bullet("The tracked docs folder (docs/tracked/) has every prompt I gave the agent, every gotcha it discovered, every fix it applied. It's the real build journal -- 17 phases documented.")
    bullet("Real gotchas discovered by the AI agent: $select causes silent zero-row returns on computed Dataverse fields. Polymorphic lookups return _parentcustomerid_value at runtime but the generated types declare parentcustomerid. The command palette crashes the app if it's rendered outside HashRouter.")
    pdf.ln(1)

    subheading("Right Column -- The Dashboard Story")
    body("\"People always want to know about the dashboard. Here's the story.\"")
    pdf.ln(1)
    bullet("Phase 6c: I told the agent to build a dashboard. It reached for Chart.js. It worked, but it was heavy.")
    bullet("Phase 8: I said 'Drop Chart.js entirely and rebuild using pure CSS and SVG.' One prompt. The agent rebuilt the entire analytics dashboard -- KPI cards, donut chart, horizontal bars, progress bars -- in 705 lines. No chart library. The donut chart is literally an SVG circle with stroke-dasharray.")
    bullet("Every visualization is clickable. Click a status row, a priority bar, an account bar -- it opens a drilldown dialog with a filtered table. All done with simple filter callbacks.")
    bullet("The Kanban board is more complex at 1,117 lines, but it's the same pattern: dnd-kit for drag-drop, TanStack Query for data, Zustand for UI state, Tailwind for styling. Standard React libraries doing standard React things.")
    pdf.ln(1)

    body("\"The point isn't that AI built it. The point is that Code Apps let you use "
         "standard tools, and standard tools are what AI agents are best at. "
         "You get a flywheel: the platform doesn't fight the tools, the tools don't fight the AI, "
         "and the AI can iterate fast because it's just writing React.\"")

    stage_direction("[Point to the stats row: 6 entities, 18 phases, 705 dashboard lines, "
                    "0 chart libraries, 1 deploy command. Let the numbers speak.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 5: Patterns That Pay Off")
    stage_direction("[Advance slide. Four pattern cards. Walk through each one briskly -- "
                    "you'll show these live in the demo.]")

    body("\"These are four patterns from the app that I think are worth calling out. "
         "None of them are exotic. They're mainstream React patterns. The point is that "
         "Code Apps let you use them.\"")
    pdf.ln(1)

    subheading("TanStack Query Cache")
    body("\"The command palette -- Ctrl+K -- searches across all six entities instantly. "
         "Zero API calls. It reads the TanStack Query cache that's already populated from "
         "your normal page loads. If you've visited the accounts list, those accounts are "
         "in the cache. The command palette just reads them. No extra Dataverse round-trips.\"")
    pdf.ln(1)

    subheading("Zustand Store (27 lines)")
    body("\"The quick-create bar at the top has colored pills -- tap 'work' and it opens "
         "the action item form pre-populated with task type 'Work', on whichever page you're on. "
         "That's a 27-line Zustand store. The pill calls store.open(), the list component watches "
         "the store, auto-opens the form dialog. No prop drilling, no context wrapping.\"")
    pdf.ln(1)

    subheading("dnd-kit + Dataverse Mutations")
    body("\"The Kanban board has cross-column drag-and-drop. Drag a card to the parking lot "
         "column and it fires a Dataverse PATCH to set tdvsp_pinned = true. Drag it back out "
         "and it unsets the flag. Within-column reorder uses arrayMove and persists to localStorage. "
         "I'll show you this live.\"")
    pdf.ln(1)

    subheading("Tailwind v4 Dark Mode")
    body("\"This is a subtle one. Media-query-based dark mode doesn't work in Code Apps "
         "because the Power Platform host controls the outer page. You need class-based dark mode. "
         "Tailwind v4 lets you declare a custom variant -- @custom-variant dark with a :where selector -- "
         "and then all your dark: utilities just work inside the iframe.\"")

    stage_direction("[Keep this slide brisk. 60-90 seconds total. "
                    "The demo will show all four patterns live.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 6: Let Me Show You")
    stage_direction("[Advance slide. Pause. Let the audience read the demo flow.]")

    body("\"Alright, enough slides. Let me show you the app.\"")
    pdf.ln(1)

    body("\"Here's what we'll walk through:\"")
    bullet("The analytics dashboard -- KPIs, the CSS donut chart, and clicking through to drilldown tables.")
    bullet("The Kanban board -- I'll drag something to the parking lot so you can see the Dataverse mutation.")
    bullet("Ctrl+K command palette -- instant search across all entities from the cache.")
    bullet("AI extraction -- I'll paste meeting notes and we'll watch Azure OpenAI create action items in real time.")
    bullet("Dark mode -- one click, class-based, works inside the Power Platform iframe.")
    bullet("Dev tools -- I'll pop open the Elements panel so you can see the runtime we just talked about.")
    bullet("Deploy -- npm run build, pac code push, done.")
    pdf.ln(2)

    stage_direction("[Switch to browser. App should already be loaded at the dashboard view. "
                    "If you need a fallback, have localhost:3001 running in a second tab.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # LIVE DEMO DETAILED BEATS
    # ════════════════════════════════════════════════════════════════
    heading("Live Demo: Detailed Beats")
    body("Below are the detailed beats for each demo segment. Times are approximate.")
    pdf.ln(2)

    subheading("Act 1: Dashboard (3 min)")
    stage_direction("[App is showing the analytics dashboard]")
    bullet("Point out the 4 KPI cards at the top. Hover over one to show the tooltip (item count, first 4 names).")
    bullet("\"This donut chart is pure SVG. No Chart.js, no D3. Just circle elements with stroke-dasharray. 705 lines total.\"")
    bullet("Click a status row (e.g., 'In Progress') to open the drilldown dialog. \"Every visualization drills down to a filtered table.\"")
    bullet("Close the drilldown. Click a priority bar to show the same drilldown with a different filter.")
    pdf.ln(1)

    subheading("Act 2: Kanban Board (4 min)")
    stage_direction("[Click 'My Board' in the sidebar]")
    bullet("Orient: 4 columns -- parking lot (green), work (dynamic), projects (purple), ideas (amber).")
    bullet("Hover over a card to show the floating toolbar (grip handle, color dots, edit pencil, car/pin icon).")
    bullet("Drag a card to the parking lot. \"That just wrote tdvsp_pinned = true to Dataverse.\" Drag it back out.")
    bullet("Click the task-type filter pills in the work column header (A/W/P/L). \"The column title, icon, and accent color all change.\"")
    bullet("Click the edit pencil on a card. \"This opens the right form dialog for the entity type -- action item, project, idea, or meeting summary.\"")
    pdf.ln(1)

    subheading("Act 3: Command Palette (2 min)")
    stage_direction("[Press Ctrl+K]")
    bullet("Type a search term that matches across entities (e.g., a company name).")
    bullet("\"Zero API calls. This reads the TanStack Query cache. If you've visited a page, those records are already here.\"")
    bullet("Arrow down and press Enter to navigate to a result. Show that it routes correctly.")
    pdf.ln(1)

    subheading("Act 4: AI Extraction -- The Showstopper (5 min)")
    stage_direction("[Navigate to Meeting Summaries. Open a meeting summary that has notes, "
                    "or create one and paste in sample meeting notes.]")
    bullet("Click 'Extract Action Items with AI'.")
    bullet("\"This sends the meeting text to Azure OpenAI. The system prompt tells it to return structured JSON -- name, priority, due date, notes.\"")
    bullet("Watch the extraction results appear. Point out that priorities are already mapped to Dataverse choice keys (High = 468510003, Medium = 468510001, Low = 468510000).")
    bullet("Bulk-create the action items. Navigate to the action items list to show them.")
    bullet("\"80 lines of code in azure-openai.ts. That's the entire integration. If the env vars aren't set, the button just doesn't appear. Graceful degradation.\"")
    pdf.ln(1)

    subheading("Act 5: Dark Mode + Dev Tools (2 min)")
    stage_direction("[Click the moon/sun icon in the sidebar footer]")
    bullet("Toggle dark mode. \"Class-based, persists to localStorage, works inside the Power Platform iframe.\"")
    bullet("Open browser dev tools (F12). Show the Elements panel.")
    bullet("\"Remember slide 2? This is the runtime. You can see the iframe nesting, the 77 registered scripts, the feature gates. Your React app is in here -- but it's a small island inside all of this.\"")
    pdf.ln(1)

    subheading("Act 6: Deploy (1 min)")
    stage_direction("[Switch to terminal -- or just narrate this if time is short]")
    code("npm run build && pac code push")
    bullet("\"That's it. Vite builds the SPA, pac code push deploys it to Power Platform. Same governance, same DLP, same security as a Canvas app. No Azure App Service, no extra hosting.\"")
    pdf.ln(1)

    subheading("Act 7: Wrap-Up (1 min)")
    body("\"So here's the story: a modern React SPA, built entirely by AI agents, "
         "running inside the Power Platform runtime with full Dataverse integration, "
         "AI-powered features, and a one-command deploy. "
         "Code Apps give your developers a path into Power Platform without giving up "
         "any of the tools or patterns they already know.\"")
    pdf.ln(2)

    divider()

    # ════════════════════════════════════════════════════════════════
    # RECOVERY PLAYS
    # ════════════════════════════════════════════════════════════════
    heading("Recovery Plays")
    body("Things that might go wrong during the live demo and how to handle them.")
    pdf.ln(1)

    subheading("Dataverse Latency")
    body("If creates or updates are slow: \"Dataverse is in the cloud -- sometimes there's a beat. "
         "In production you'd see this as a loading spinner. The TanStack Query invalidation "
         "will pick up the change as soon as it lands.\"")

    subheading("AI Extraction Fails")
    body("If Azure OpenAI times out or returns unexpected results: \"This is a live API call -- "
         "sometimes the model takes a moment. The error handling here shows a toast with the "
         "specific error. In production you'd add retry logic. The point is the integration "
         "pattern, not the uptime of my demo key.\"")

    subheading("\"Is this Canvas or Model-Driven?\"")
    body("\"Neither, and both. Code Apps are a third option. You write React, you deploy to "
         "Power Platform, you get the same Dataverse backend, the same governance, the same "
         "security model. The runtime even classifies it as ClassicCanvasApp internally. "
         "But your source code is standard React/TypeScript.\"")

    subheading("\"Can I use this with Copilot Studio?\"")
    body("\"Yes. The runtime has a Copilot sidecar built in -- we saw 185 references to it "
         "in the shell HTML. And the AI extraction feature in this demo uses Azure OpenAI directly. "
         "Both patterns work.\"")

    subheading("\"How do I get started?\"")
    body("\"My colleague covered that in the intro. pac code init, pac code add-data-source, "
         "npm run dev. The getting-started story is straightforward. What I showed you today "
         "is what happens when you push past hello world.\"")

    pdf_path = os.path.join(OUT_DIR, "code-apps-under-the-hood-talk-track.pdf")
    pdf.output(pdf_path)
    print(f"PDF saved: {pdf_path}")
    return pdf_path


# ── Main ────────────────────────────────────────────────────────────
if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print("Done!")
